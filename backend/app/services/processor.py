"""
DocumentProcessor — Extractor Premium de Texto Multi-Formato v3
===============================================================
Mejoras v3:
  - IMÁGENES: OCR real con Gemini Vision (lee texto, tablas, gráficos)
  - PDF escaneado: detecta páginas sin texto y aplica Gemini Vision
  - PPTX: extrae texto de presentaciones PowerPoint
  - Excel: soporta múltiples hojas con nombre de hoja
  - Compresion agresiva de texto: elimina redundancias, páginas vacías
  - Extraccion por páginas con filtro de calidad
  - Estadísticas de extracción por documento
"""

import fitz          # PyMuPDF
from docx import Document
import pandas as pd
import io
import os
import base64
import logging
from PIL import Image
from typing import Dict, List, Tuple, Optional
import re

logger = logging.getLogger(__name__)

# ── Configuración Gemini Vision (para imágenes y PDFs escaneados) ────────────
_gemini_vision = None

def _get_gemini_vision():
    """Inicializa modelo Gemini Vision de forma lazy."""
    global _gemini_vision
    if _gemini_vision is None:
        api_key = os.environ.get("GOOGLE_API_KEY", "")
        if api_key:
            try:
                import google.generativeai as genai
                genai.configure(api_key=api_key)
                _gemini_vision = genai.GenerativeModel(
                    model_name="gemini-2.0-flash",
                    generation_config=genai.GenerationConfig(
                        temperature=0.0,
                        max_output_tokens=8192,
                    )
                )
                logger.info("[VISION] Gemini Vision listo (gemini-2.0-flash)")
            except Exception as e:
                logger.warning(f"[VISION] Gemini Vision no disponible: {e}")
    return _gemini_vision


def _ocr_con_gemini(image_bytes: bytes, contexto: str = "") -> str:
    """
    Usa Gemini Vision para extraer texto completo de una imagen.
    Incluye tablas, gráficos, texto manuscrito, etc.
    """
    model = _get_gemini_vision()
    if model is None:
        return "[OCR no disponible: GOOGLE_API_KEY no configurada]"

    try:
        import google.generativeai as genai

        # Detectar formato de imagen
        img = Image.open(io.BytesIO(image_bytes))
        fmt = img.format or "PNG"
        mime = f"image/{fmt.lower()}"
        if mime == "image/jpg":
            mime = "image/jpeg"

        prompt_parts = [
            genai.protos.Part(
                inline_data=genai.protos.Blob(
                    mime_type=mime,
                    data=image_bytes
                )
            ),
            genai.protos.Part(text=(
                "Eres un extractor de información. Analiza esta imagen (que puede ser una "
                "página de documento, tabla, gráfico, foto, captura de pantalla, etc.) y:\n"
                "1. Extrae TODO el texto visible, conservando la estructura\n"
                "2. Si hay tablas, conviértelas a formato texto plano con columnas separadas por ' | '\n"
                "3. Si hay gráficos/charts, describe los datos principales que muestra\n"
                "4. Si hay fórmulas, escríbelas en texto\n"
                "5. Si hay texto manuscrito, transcríbelo\n"
                f"{f'Contexto del documento: {contexto}' if contexto else ''}\n"
                "Responde SOLO con el contenido extraído, sin explicaciones adicionales."
            ))
        ]

        response = model.generate_content(prompt_parts)
        texto = response.text.strip()
        logger.info(f"[VISION OCR] Extraído: {len(texto):,} chars")
        return texto

    except Exception as e:
        logger.error(f"[VISION OCR] Error: {e}")
        return f"[Error en OCR de imagen: {e}]"


class DocumentProcessor:

    # ── Limpieza de Texto ───────────────────────────────────────────────────────

    @staticmethod
    def clean_text(text: str) -> str:
        """Limpieza premium: elimina espacios excesivos y caracteres basura."""
        text = re.sub(r'\n{3,}', '\n\n', text)
        lines = text.split('\n')
        clean_lines = []
        for line in lines:
            stripped = line.strip()
            if re.match(r'^\d{1,4}$', stripped):
                continue
            if len(stripped) < 3 and stripped not in ['', '\n']:
                continue
            clean_lines.append(line)
        text = '\n'.join(clean_lines)
        text = re.sub(r'[ \t]{2,}', ' ', text)
        return text.strip()

    @staticmethod
    def compress_text(text: str) -> str:
        """
        Compresión inteligente: elimina contenido redundante para maximizar
        la cantidad de información útil por token enviado a la IA.
        """
        lines = text.split('\n')
        seen_lines = {}
        compressed = []
        for line in lines:
            key = line.strip().lower()
            if not key:
                compressed.append(line)
                continue
            count = seen_lines.get(key, 0)
            if count < 2:
                compressed.append(line)
                seen_lines[key] = count + 1
        return '\n'.join(compressed)

    # ── Extracción por Formato ──────────────────────────────────────────────────

    def _extract_pdf(self, file_content: bytes, filename: str) -> Tuple[str, dict]:
        """
        Extrae texto de PDF con filtro de calidad por página.
        Para páginas escaneadas (sin texto), usa Gemini Vision OCR.
        """
        parts = []
        stats = {
            "pages_total": 0,
            "pages_extracted": 0,
            "pages_skipped": 0,
            "pages_ocr": 0
        }

        with fitz.open(stream=file_content, filetype="pdf") as doc:
            stats["pages_total"] = len(doc)
            for i, page in enumerate(doc):
                page_text = page.get_text("text")
                cleaned = self.clean_text(page_text)
                words = len(cleaned.split())

                if words >= 10:
                    # Página con texto normal
                    parts.append(f"[PAG.{i+1}] {cleaned}")
                    stats["pages_extracted"] += 1
                else:
                    # Página escaneada o con muy poco texto → intentar OCR con Gemini Vision
                    model = _get_gemini_vision()
                    if model is not None:
                        try:
                            # Renderizar la página como imagen
                            mat = fitz.Matrix(2, 2)  # 2x zoom para mejor calidad
                            pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
                            img_bytes = pix.tobytes("png")
                            ocr_text = _ocr_con_gemini(
                                img_bytes,
                                contexto=f"Página {i+1} del PDF '{filename}'"
                            )
                            if ocr_text and len(ocr_text.strip()) > 20:
                                parts.append(f"[PAG.{i+1}][OCR] {ocr_text}")
                                stats["pages_extracted"] += 1
                                stats["pages_ocr"] += 1
                            else:
                                stats["pages_skipped"] += 1
                        except Exception as e:
                            logger.warning(f"[PDF OCR] Página {i+1}: {e}")
                            stats["pages_skipped"] += 1
                    else:
                        stats["pages_skipped"] += 1

        return '\n'.join(parts), stats

    def _extract_docx(self, file_content: bytes) -> str:
        """Extrae texto de Word con estructura de párrafos y tablas."""
        doc = Document(io.BytesIO(file_content))
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                if para.style.name.startswith('Heading'):
                    parts.append(f"\n## {text}")
                else:
                    parts.append(text)
        # Tablas de Word
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    rows.append(' | '.join(cells))
            if rows:
                parts.append('[TABLA]\n' + '\n'.join(rows))
        return '\n'.join(parts)

    def _extract_excel(self, file_content: bytes, ext: str) -> str:
        """
        Extrae datos tabulares de Excel/CSV.
        Soporta múltiples hojas con nombre de hoja.
        """
        try:
            if ext == 'csv':
                df = pd.read_csv(io.BytesIO(file_content), encoding='utf-8', errors='ignore')
                df = df.dropna(axis=1, how='all').fillna('')
                if len(df) > 300:
                    return (f"[TABLA CSV — mostrando 300 de {len(df)} filas]\n"
                            f"{df.head(300).to_string(index=False)}")
                return f"[TABLA CSV — {len(df)} filas x {len(df.columns)} cols]\n{df.to_string(index=False)}"

            else:
                # Excel con múltiples hojas
                xl = pd.ExcelFile(io.BytesIO(file_content))
                all_sheets = []
                for sheet_name in xl.sheet_names:
                    try:
                        df = xl.parse(sheet_name)
                        df = df.dropna(axis=1, how='all').fillna('')
                        if df.empty:
                            continue
                        # Limitar a 200 filas por hoja
                        if len(df) > 200:
                            texto = f"[HOJA: '{sheet_name}' — mostrando 200 de {len(df)} filas]\n{df.head(200).to_string(index=False)}"
                        else:
                            texto = f"[HOJA: '{sheet_name}' — {len(df)} filas x {len(df.columns)} cols]\n{df.to_string(index=False)}"
                        all_sheets.append(texto)
                    except Exception as e:
                        all_sheets.append(f"[HOJA: '{sheet_name}' — Error: {e}]")

                return '\n\n'.join(all_sheets) if all_sheets else "[Excel sin datos]"

        except Exception as e:
            return f"[Error leyendo tabla: {e}]"

    def _extract_image(self, file_content: bytes, filename: str) -> str:
        """
        Extrae información de imágenes usando Gemini Vision OCR.
        Lee texto, tablas, gráficos, capturas de pantalla, etc.
        """
        try:
            img = Image.open(io.BytesIO(file_content))
            info = f"Formato: {img.format}, Tamaño: {img.size[0]}x{img.size[1]}px, Modo: {img.mode}"

            # OCR con Gemini Vision
            model = _get_gemini_vision()
            if model is not None:
                ocr_text = _ocr_con_gemini(
                    file_content,
                    contexto=f"Imagen '{filename}'"
                )
                return f"[IMAGEN] {info}\n[CONTENIDO EXTRAÍDO POR IA]\n{ocr_text}"
            else:
                return f"[IMAGEN] {info}\n[OCR no disponible — configura GOOGLE_API_KEY]"

        except Exception as e:
            return f"[Error procesando imagen: {e}]"

    def _extract_pptx(self, file_content: bytes) -> str:
        """Extrae texto de presentaciones PowerPoint."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_content))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_parts = [f"[DIAPOSITIVA {i}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                if len(slide_parts) > 1:  # Tiene contenido además del encabezado
                    parts.append('\n'.join(slide_parts))
            return '\n\n'.join(parts) if parts else "[Presentación sin texto extraíble]"
        except ImportError:
            return "[Instala python-pptx para soporte PPTX: pip install python-pptx]"
        except Exception as e:
            return f"[Error leyendo PPTX: {e}]"

    # ── Método Principal ────────────────────────────────────────────────────────

    def extract_text_with_metadata(self, file_content: bytes, filename: str) -> str:
        """
        Extrae texto de múltiples formatos con limpieza y compresión premium.
        Soporta: PDF, Word, Excel, CSV, TXT, MD, imágenes (con OCR IA), PPTX.
        """
        ext = filename.split('.')[-1].lower() if '.' in filename else ''
        header = f"\n{'='*60}\n[[ DOCUMENTO: {filename.upper()} ]]\n{'='*60}\n"

        try:
            if ext == 'pdf':
                raw_text, stats = self._extract_pdf(file_content, filename)
                compressed = self.compress_text(raw_text)
                ocr_nota = f" | {stats['pages_ocr']} págs. con OCR" if stats.get('pages_ocr', 0) > 0 else ""
                info = (
                    f"[INFO: {stats['pages_extracted']}/{stats['pages_total']} págs. extraídas"
                    f" | {stats['pages_skipped']} págs. omitidas{ocr_nota}"
                    f" | {len(compressed):,} chars]\n"
                )
                return header + info + compressed

            elif ext in ['docx', 'doc']:
                raw_text = self._extract_docx(file_content)
                compressed = self.compress_text(self.clean_text(raw_text))
                return header + f"[INFO: {len(compressed):,} chars extraídos de Word]\n" + compressed

            elif ext in ['xlsx', 'xls', 'csv']:
                table_text = self._extract_excel(file_content, ext)
                return header + table_text

            elif ext in ['txt', 'md']:
                raw = file_content.decode('utf-8', errors='ignore')
                compressed = self.compress_text(self.clean_text(raw))
                return header + f"[INFO: {len(compressed):,} chars]\n" + compressed

            elif ext in ['jpg', 'jpeg', 'png', 'webp', 'bmp', 'tiff', 'gif']:
                image_text = self._extract_image(file_content, filename)
                return header + image_text

            elif ext in ['pptx', 'ppt']:
                pptx_text = self._extract_pptx(file_content)
                compressed = self.compress_text(self.clean_text(pptx_text))
                return header + f"[INFO: {len(compressed):,} chars extraídos de PowerPoint]\n" + compressed

            else:
                # Intentar leer como texto plano (último recurso)
                try:
                    raw = file_content.decode('utf-8', errors='ignore')
                    if raw.strip():
                        compressed = self.compress_text(self.clean_text(raw))
                        return header + f"[TEXTO PLANO — {len(compressed):,} chars]\n" + compressed
                except Exception:
                    pass
                return header + f"[Formato '{ext}' no soportado. Tamaño: {len(file_content):,} bytes]\n"

        except Exception as e:
            logger.error(f"[PROCESSOR] Error extrayendo {filename}: {e}")
            return header + f"[ERROR DE EXTRACCIÓN: {str(e)}]\n"

    def extract_documents_individually(self, documents: Dict[str, bytes]) -> List[Dict]:
        """
        Extrae texto de cada documento por separado.
        Retorna una lista de dicts con metadata por documento:
          { name, text, char_count, token_estimate }
        """
        result = []
        for doc_name, content in documents.items():
            logger.info(f"[PROCESSOR] Procesando: '{doc_name}' ({len(content):,} bytes)")
            text = self.extract_text_with_metadata(content, doc_name)
            char_count = len(text)
            result.append({
                "name": doc_name,
                "text": text,
                "char_count": char_count,
                "token_estimate": char_count // 4,
            })
            logger.info(
                f"[PROCESSOR] '{doc_name}': {char_count:,} chars "
                f"(~{char_count // 4:,} tokens)"
            )
        return result

    @staticmethod
    def validate_integrity(documents: Dict[str, bytes]) -> List[str]:
        """
        Checklist de integridad opcional (no bloqueante).
        """
        required = [
            "PEI", "MANUAL DE CONVIVENCIA", "PMI", "POA",
            "PFI", "SIEE", "LECTURA DE CONTEXTO"
        ]
        present_docs = [name.upper() for name in documents.keys()]
        return [doc for doc in required if doc not in present_docs]

    def prepare_context_for_ai(self, documents: Dict[str, bytes]) -> str:
        """
        Prepara el contexto completo concatenado para la IA.
        """
        parts = self.extract_documents_individually(documents)
        total_chars = sum(p["char_count"] for p in parts)
        logger.info(
            f"[PROCESSOR] Total contexto: {total_chars:,} chars "
            f"({len(parts)} docs)"
        )
        return "\n\n".join(p["text"] for p in parts)


processor = DocumentProcessor()
